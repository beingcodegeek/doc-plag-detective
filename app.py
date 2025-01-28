import os
import re
import zipfile
import hashlib
from flask import Flask, request, jsonify
from flask_cors import CORS
import fitz  # PyMuPDF for PDF processing
from docx import Document
from pptx import Presentation
from sklearn.feature_extraction.text import CountVectorizer
import numpy as np

app = Flask(__name__)
CORS(app)
TEMP_DIR = "temp_files"
os.makedirs(TEMP_DIR, exist_ok=True)

def preprocess_text(text):
    """Lowercase and remove punctuation."""
    if not text:
        return ""
    text = text.lower()
    text = re.sub(r"[\W_]+", " ", text)
    return text

def extract_text_from_pdf(file):
    """Extract text from PDF using PyMuPDF."""
    text = ""
    pdf_document = fitz.open(stream=file.read(), filetype="pdf")
    for page in pdf_document:
        text += page.get_text()
    return text.strip()

def extract_text_from_docx(file):
    """Extract text from DOCX files."""
    doc = Document(file)
    return " ".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(file):
    """Extract text from PPTX files."""
    presentation = Presentation(file)
    return " ".join([shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text")])

def process_file(file, filename):
    """Process supported file types."""
    if filename.endswith(".pdf"):
        return extract_text_from_pdf(file)
    elif filename.endswith(".docx"):
        return extract_text_from_docx(file)
    elif filename.endswith(".pptx"):
        return extract_text_from_pptx(file)
    return None

def cosine_similarity(text1, text2):
    """Calculate the cosine similarity between two texts."""
    vectorizer = CountVectorizer()
    vectors = vectorizer.fit_transform([text1, text2])
    cosine_sim = (vectors * vectors.T).toarray()
    return cosine_sim[0, 1]

def jaccard_similarity(text1, text2):
    """Calculate Jaccard similarity between two texts."""
    set1 = set(text1.split())
    set2 = set(text2.split())
    intersection = len(set1.intersection(set2))
    union = len(set1.union(set2))
    return intersection / union if union != 0 else 0

@app.route('/upload', methods=['POST'])
def upload_files():
    file_details = request.json.get('fileDetails', [])
    if 'fileDetails' not in request.json:
        return jsonify({"success": False, "message": "Invalid request data"}), 400

    file_paths = []
    files = {}
    for file in file_details:
        student_id = file.get('studentId')
        file_url = file.get('fileUrl')

        if not file_url or not student_id:
            continue

        response = requests.get(file_url)
        if response.status_code == 200:
            file_name = f"{os.path.basename(file_url)}"
            file_path = os.path.join(TEMP_DIR, file_name)
            files[file_name] = student_id
            with open(file_path, 'wb') as f:
                f.write(response.content)
            file_paths.append(file_path)

        else:
            print(f"Failed to download file from: {file_url}")
  
    uploaded_files = []
    assignments = {}
    for file in file_paths:
        uploaded_files.append(file)

    for file in uploaded_files:
        with open(file, "rb") as f:
            text = process_file(f, file)
            if text:
                assignments[os.path.basename(file)] = preprocess_text(text)
            else:
                text = process_file(file, file)
    
    # Calculate similarities
    results = []
    assignment_names = list(assignments.keys())
    for i in range(len(assignment_names)):
        for j in range(i + 1, len(assignment_names)):
            name1 = assignment_names[i]
            name2 = assignment_names[j]

            # Calculate Cosine Similarity
            cosine_sim = cosine_similarity(assignments[name1], assignments[name2])

            # Calculate Jaccard Similarity
            jaccard_sim = jaccard_similarity(assignments[name1], assignments[name2])

            # Combine both similarities into a final combined score
            combined_sim = (cosine_sim + jaccard_sim) / 2

            results.append({
                "Assignment 1": name1,
                "Assignment 2": name2,
                "Cosine Similarity (%)": round(cosine_sim * 100, 2),
                "Jaccard Similarity (%)": round(jaccard_sim * 100, 2),
                "Combined Similarity (%)": round(combined_sim * 100, 2)
            })

    return jsonify({"results": results})

if __name__ == '__main__':
    app.run(port=8081)
